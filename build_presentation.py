#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Создаёт презентацию PPTX по теме магистерской диссертации.
Слайды:
  1. Титул
  2. Актуальность
  3. Объект / Предмет / Цель
  4. Задачи исследования
  5. Анализ существующих методов
  6. Датасет и данные
  7. Архитектура решения
  8. Предложенное решение
  9. Результаты: метрики
 10. Результаты: графики ROC/PR
 11. Результаты: динамика score
 12. Выводы и направления развития
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path("Presentation_Anomaly_Detection.pptx")

# ─── цветовая палитра ────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1F, 0x35, 0x64)   # тёмно-синий
C_ACCENT = RGBColor(0x2E, 0x75, 0xB6)   # синий акцент
C_LIGHT  = RGBColor(0xD9, 0xEA, 0xF7)   # голубой фон
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x40, 0x40, 0x40)
C_GREEN  = RGBColor(0x54, 0x82, 0x35)
C_ORANGE = RGBColor(0xC6, 0x59, 0x11)

FIG = Path("figures")

# ─── вспомогательные функции ─────────────────────────────────────────────────

def blank_slide(prs):
    blank = prs.slide_layouts[6]          # blank layout
    return prs.slides.add_slide(blank)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, *,
             size=20, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, lines, l, t, w, h, *,
                   title=None, title_size=18, body_size=15,
                   title_color=C_DARK, body_color=C_GRAY,
                   bullet="•  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(title_size)
        r.font.bold  = True
        r.font.color.rgb = title_color
        first = False
    for line in lines:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        r = p.add_run()
        r.text = bullet + line
        r.font.size  = Pt(body_size)
        r.font.color.rgb = body_color
        p.space_before = Pt(4)


def add_image(slide, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists():
        return
    if h:
        slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))


def slide_header(slide, title_text, subtitle=None):
    """Синяя полоса сверху + заголовок."""
    add_rect(slide, 0, 0, 13.33, 1.15, C_DARK)
    add_text(slide, title_text, 0.35, 0.12, 12.6, 0.9,
             size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.88, 12.6, 0.35,
                 size=13, color=C_LIGHT, align=PP_ALIGN.LEFT)


# ─── создаём презентацию ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_DARK)

# Декоративная полоса
add_rect(sl, 0, 5.6, 13.33, 0.08, C_ACCENT)

add_text(sl, "МАГИСТЕРСКАЯ ДИССЕРТАЦИЯ", 1.0, 0.55, 11.33, 0.6,
         size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_text(sl,
         "Разработка системы детекции аномального\n"
         "сетевого трафика с использованием\n"
         "методов машинного обучения",
         0.8, 1.2, 11.73, 2.4,
         size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Направление 10.04.01 — Информационная безопасность",
         1.0, 3.8, 11.33, 0.5, size=15, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, 4.0, 4.4, 5.33, 0.06, C_ACCENT)

add_text(sl, "2026", 1.0, 4.7, 11.33, 0.5,
         size=14, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — АКТУАЛЬНОСТЬ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Актуальность исследования")

# Левая карточка
add_rect(sl, 0.3, 1.35, 4.0, 5.65, C_LIGHT)
add_text(sl, "Вызовы", 0.5, 1.45, 3.6, 0.45, size=15, bold=True, color=C_DARK)
add_bullet_box(sl,
    ["Рост зашифрованного трафика (TLS повсеместно)",
     "Сигнатурные IDS не покрывают неизвестные атаки",
     "Разметка инцидентов дорогостояща и редка",
     "Операторы страдают от «усталости алертов»",
     "Необходимость локального развёртывания (закрытые сети)"],
    0.5, 1.95, 3.6, 4.8, body_size=13, bullet="▸  ")

# Центральная карточка
add_rect(sl, 4.65, 1.35, 4.0, 5.65, RGBColor(0xE2, 0xF0, 0xD9))
add_text(sl, "Тенденции", 4.85, 1.45, 3.6, 0.45, size=15, bold=True, color=C_GREEN)
add_bullet_box(sl,
    ["Переход к анализу поведения потоков вместо payload",
     "ML без учителя как основа NDR-систем",
     "Ансамблевые детекторы повышают устойчивость",
     "Объяснимость алертов — ключевое требование SOC",
     "Prometheus / Grafana как стандарт наблюдаемости"],
    4.85, 1.95, 3.6, 4.8, body_size=13, bullet="▸  ")

# Правая карточка
add_rect(sl, 9.0, 1.35, 4.03, 5.65, RGBColor(0xFC, 0xE4, 0xD6))
add_text(sl, "Разрыв в исследованиях", 9.2, 1.45, 3.6, 0.45,
         size=15, bold=True, color=C_ORANGE)
add_bullet_box(sl,
    ["Мало работ с воспроизводимым end-to-end конвейером",
     "Порог и FPR редко рассматриваются как требования",
     "Нет единого открытого прototипа с мониторингом",
     "Синтетические сценарии недостаточно используются для проверки"],
    9.2, 1.95, 3.6, 4.8, body_size=13, bullet="▸  ")

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — ОБЪЕКТ / ПРЕДМЕТ / ЦЕЛЬ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Объект, предмет и цель исследования")

blocks = [
    (C_LIGHT,  C_ACCENT, "Объект исследования",
     "Сетевой трафик в компьютерных сетях и процессы его мониторинга"),
    (RGBColor(0xE2, 0xF0, 0xD9), C_GREEN, "Предмет исследования",
     "Методы и программные средства детекции аномального трафика "
     "на основе статистических признаков и ML-моделей, "
     "а также критерии качества таких систем — точность, задержка, устойчивость"),
    (RGBColor(0xFC, 0xE4, 0xD6), C_ORANGE, "Цель исследования",
     "Разработать прototип системы детекции аномального сетевого трафика, "
     "использующей методы машинного обучения и признаки потоков/окон, "
     "и оценить её качество на воспроизводимых сценариях"),
]
y = 1.35
for bg, accent, label, body in blocks:
    add_rect(sl, 0.3, y, 12.73, 1.65, bg)
    add_rect(sl, 0.3, y, 0.1, 1.65, accent)
    add_text(sl, label, 0.6, y + 0.12, 12.0, 0.4,
             size=16, bold=True, color=accent)
    add_text(sl, body, 0.6, y + 0.55, 12.0, 0.95,
             size=13.5, color=C_GRAY)
    y += 1.85

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — ЗАДАЧИ ИССЛЕДОВАНИЯ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Задачи исследования")

tasks = [
    "Обзор и сравнение подходов к детекции сетевых аномалий (сигнатурные, статистические, ML)",
    "Обоснование выбора источника данных: пакетная телеметрия и Parquet-датасет",
    "Формирование признакового пространства: потоки/окна без анализа payload",
    "Реализация сбора и хранения с партиционированием date/hour",
    "Обучение ансамбля (IForest + AutoEncoder) и выбор порогов T_low / T_high",
    "Определение метрик качества: Precision, Recall, F1, ROC-AUC, PR-AUC, FPR",
    "Проектирование архитектуры: сбор → признаки → ML → углублённая проверка → мониторинг",
]
cols = [tasks[:4], tasks[4:]]
x_positions = [0.35, 6.85]
for col, x in zip(cols, x_positions):
    for idx, task in enumerate(col):
        y = 1.4 + idx * 1.45
        add_rect(sl, x, y, 6.2, 1.3, C_LIGHT)
        add_rect(sl, x, y, 0.45, 1.3, C_ACCENT)
        add_text(sl, str(tasks.index(task) + 1), x + 0.05, y + 0.35, 0.35, 0.5,
                 size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, task, x + 0.6, y + 0.15, 5.45, 1.0, size=12.5, color=C_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — АНАЛИЗ СУЩЕСТВУЮЩИХ МЕТОДОВ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Анализ существующих методов детекции")

add_image(sl, FIG / "chapter1" / "figure_1_2_detection_methods_v2.png",
          0.3, 1.25, 5.5)

# Таблица сравнения справа
rows = [
    ("Сигнатурные IDS/IPS",   "Высокая",   "Новые атаки",      "Задержка, FPR"),
    ("Статистические методы", "Средняя",   "Дрейф фона",       "Простота"),
    ("ML без учителя",        "Высокая",   "FPR, объяснимость","Адаптивность"),
    ("NDR-платформы",         "Очень выс.","Прозрачность",     "Поведение"),
]
headers = ["Метод", "Объясн.", "Слабость", "Сильная сторона"]
col_w = [2.35, 1.0, 1.65, 1.65]
col_x = [6.1, 8.45, 9.45, 11.1]
y0 = 1.3

# заголовки таблицы
add_rect(sl, 6.1, y0, 6.85, 0.4, C_DARK)
for i, h in enumerate(headers):
    add_text(sl, h, col_x[i] + 0.05, y0 + 0.05, col_w[i] - 0.1, 0.35,
             size=11, bold=True, color=C_WHITE)

for ri, row in enumerate(rows):
    bg = C_LIGHT if ri % 2 == 0 else C_WHITE
    add_rect(sl, 6.1, y0 + 0.4 + ri * 0.55, 6.85, 0.55, bg)
    for ci, cell in enumerate(row):
        clr = C_ACCENT if ci == 0 else C_GRAY
        add_text(sl, cell, col_x[ci] + 0.05, y0 + 0.45 + ri * 0.55,
                 col_w[ci] - 0.1, 0.48, size=11, color=clr,
                 bold=(ci == 0))

# Вывод
add_rect(sl, 6.1, y0 + 0.4 + len(rows) * 0.55 + 0.15, 6.85, 0.85,
         RGBColor(0xE2, 0xF0, 0xD9))
add_text(sl, "Вывод: гибридный подход (ML + правила) обеспечивает "
             "высокий recall при контролируемом FPR и объяснимые алерты",
         6.25, y0 + 0.55 + len(rows) * 0.55 + 0.15, 6.5, 0.7,
         size=12, color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — ДАТАСЕТ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Данные и датасет")

# Левая колонка — источники
add_rect(sl, 0.3, 1.3, 5.8, 5.7, C_LIGHT)
add_text(sl, "Открытые источники данных", 0.5, 1.4, 5.4, 0.45,
         size=15, bold=True, color=C_DARK)

datasets = [
    ("CICIDS2017/2018", "Canadian Institute for Cybersecurity.\nTCP/UDP потоки с метками атак:\nBotNet, DDoS, PortScan, Brute Force"),
    ("UNSW-NB15",       "University of NSW.\n49 признаков потоков, 9 классов атак.\nИспользуется в задачах ML-детекции"),
    ("Synthetic data",  "Генератор synthetic_generator.py.\nBurst (×5 интенсивность) и Port Scan.\nВоспроизводимые сценарии, известная разметка"),
]
y = 1.95
for name, desc in datasets:
    add_rect(sl, 0.4, y, 5.6, 1.45, C_WHITE)
    add_text(sl, name, 0.55, y + 0.08, 5.2, 0.35,
             size=13, bold=True, color=C_ACCENT)
    add_text(sl, desc, 0.55, y + 0.42, 5.2, 0.95, size=11, color=C_GRAY)
    y += 1.6

# Правая колонка — формат и предобработка
add_rect(sl, 6.55, 1.3, 6.48, 5.7, RGBColor(0xF8, 0xF8, 0xF8))
add_text(sl, "Формат хранения и предобработка", 6.75, 1.4, 6.1, 0.45,
         size=15, bold=True, color=C_DARK)

steps = [
    ("Захват / загрузка", "Scapy/Npcap или CSV-импорт CICIDS → Parquet"),
    ("Партиционирование", "Разбивка по date/hour для эффективной обработки"),
    ("Группировка",       "flow_key = src_ip:sport — dst_ip:dport — protocol"),
    ("Признаки (16-мер.)", "15 числовых + protocol_enc; без raw payload"),
    ("Нормализация",      "StandardScaler, сохранён внутри ensemble.joblib"),
    ("Метки",             "is_anomaly — известны из датасета или генератора"),
]
y2 = 1.95
for label, desc in steps:
    add_rect(sl, 6.65, y2, 6.18, 0.82, C_WHITE)
    add_rect(sl, 6.65, y2, 0.08, 0.82, C_ACCENT)
    add_text(sl, label, 6.85, y2 + 0.04, 3.0, 0.35, size=12, bold=True, color=C_ACCENT)
    add_text(sl, desc,  6.85, y2 + 0.38, 5.85, 0.38, size=11, color=C_GRAY)
    y2 += 0.92

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — АРХИТЕКТУРА
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Архитектура системы")

add_image(sl, FIG / "chapter2" / "figure_2_1_traffic_pipeline.png",
          0.3, 1.25, 12.73, 3.0)

# Описание модулей внизу
modules = [
    ("Коллектор",   "Scapy/Npcap\nParquet, date/hour"),
    ("Признаки",    "16 признаков\nflow_key агрег."),
    ("ML-слой",     "IForest + AE\nAnomaly score"),
    ("Проверка",    "5 правил\ndeep_inspection"),
    ("Мониторинг",  "Prometheus\nGrafana"),
]
x = 0.25
for mod, desc in modules:
    add_rect(sl, x, 4.45, 2.4, 1.6, C_LIGHT)
    add_rect(sl, x, 4.45, 2.4, 0.38, C_ACCENT)
    add_text(sl, mod, x + 0.12, 4.48, 2.15, 0.35, size=13, bold=True, color=C_WHITE)
    add_text(sl, desc, x + 0.15, 4.9, 2.1, 0.95, size=11.5, color=C_GRAY)
    x += 2.6

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — ПРЕДЛОЖЕННОЕ РЕШЕНИЕ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Предложенное решение")

# Левая — ансамбль
add_image(sl, FIG / "chapter2" / "figure_2_4_ml_ensemble.png",
          0.3, 1.25, 6.0)

# Правая — ключевые решения
items = [
    (C_ACCENT, "Безучительский ансамбль",
     "IForest (n=200) + AutoEncoder обучаются только\n"
     "на нормальном трафике. Итоговый score = среднее\n"
     "min-max нормированных оценок обоих алгоритмов."),
    (C_GREEN, "Двухпороговая схема",
     "T_low = 90-й процентиль → suspicious\n"
     "T_high = 99-й процентиль → malicious\n"
     "(только при наличии rule:* подтверждения)"),
    (C_ORANGE, "Углублённая проверка",
     "5 детерминированных правил (TCP SYN low ports,\n"
     "high burst, high entropy, compressed timeline,\n"
     "high ML score) — объяснимые алерты для SOC"),
]
y = 1.3
for color, title, body in items:
    add_rect(sl, 6.65, y, 6.38, 1.72, C_LIGHT)
    add_rect(sl, 6.65, y, 0.12, 1.72, color)
    add_text(sl, title, 6.9, y + 0.12, 5.9, 0.38,
             size=14, bold=True, color=color)
    add_text(sl, body,  6.9, y + 0.55, 5.9, 1.0, size=12, color=C_GRAY)
    y += 1.9

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — РЕЗУЛЬТАТЫ: МЕТРИКИ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Результаты: метрики качества (burst-сценарий)")

# Большие цифры
metrics = [
    ("0.970", "Recall",    C_GREEN,  "Модель обнаруживает\n97% аномалий"),
    ("0.955", "ROC-AUC",   C_ACCENT, "Высокое качество\nранжирования"),
    ("0.856", "PR-AUC",    C_ACCENT, "Устойчиво при\nдисбалансе классов"),
    ("0.641", "F1-Score",  C_DARK,   "Баланс точности\nи полноты"),
    ("0.538", "FPR",       C_ORANGE, "Требует калибровки\nпорогов"),
    ("0.478", "Precision", C_DARK,   "Улучшится после\nнастройки порога"),
]
x = 0.25
for val, name, color, note in metrics:
    add_rect(sl, x, 1.3, 2.05, 2.5, C_LIGHT)
    add_text(sl, val, x + 0.1, 1.4, 1.85, 1.0,
             size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, name, x + 0.1, 2.35, 1.85, 0.45,
             size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, note, x + 0.1, 2.82, 1.85, 0.85,
             size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
    x += 2.2

# Вывод под цифрами
add_rect(sl, 0.3, 3.95, 12.73, 0.9, RGBColor(0xE2, 0xF0, 0xD9))
add_text(sl,
         "Высокий ROC-AUC (0.955) и recall (0.97) подтверждают эффективность безучительского ансамбля. "
         "FPR=0.54 объясним особенностями синтетических данных и калибруется повышением квантиля порога до 0.995.",
         0.5, 4.0, 12.3, 0.8, size=13, color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — ROC/PR КРИВЫЕ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Результаты: ROC- и PR-кривые")

add_image(sl, FIG / "chapter3" / "figure_3_6_roc_pr_curves.png",
          0.3, 1.25, 8.3)

add_rect(sl, 8.85, 1.3, 4.18, 5.65, C_LIGHT)
add_text(sl, "Интерпретация", 9.05, 1.42, 3.8, 0.42,
         size=15, bold=True, color=C_DARK)
add_bullet_box(sl, [
    "Рабочая точка: recall=0.97, FPR=0.54",
    "Смещение вправо по ROC снижает FPR при умеренной потере recall",
    "PR-AUC=0.856 — хорошее качество при редком классе аномалий",
    "Разрыв между ROC-AUC и PR-AUC незначителен: классы относительно сбалансированы в синтетике",
    "Ансамбль (hybrid) незначительно превосходит IForest solo на burst",
], 9.05, 1.9, 3.8, 4.8, body_size=12)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 11 — ДИНАМИКА SCORE
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Результаты: динамика оценки аномальности")

add_image(sl, FIG / "chapter3" / "figure_3_8_score_timeline.png",
          0.3, 1.25, 8.3)

add_rect(sl, 8.85, 1.3, 4.18, 5.65, C_LIGHT)
add_text(sl, "Ключевые наблюдения", 9.05, 1.42, 3.8, 0.42,
         size=15, bold=True, color=C_DARK)
add_bullet_box(sl, [
    "Score резко возрастает в интервале 30–40 с (burst-аномалия)",
    "Вне интервала атаки score стабильно низкий",
    "T_low=0.611 и T_high=0.858 корректно разделяют режимы",
    "Система обнаруживает начало аномалии в первом же фрейме интервала",
    "Метрика ndr_suspicious_rate в Prometheus отражает тот же паттерн в реальном времени",
], 9.05, 1.9, 3.8, 4.8, body_size=12)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — ВЫВОДЫ
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_WHITE)
slide_header(sl, "Выводы и направления развития")

conclusions = [
    (C_GREEN,  "Реализован воспроизводимый end-to-end конвейер",
               "Scapy → Parquet → признаки → ML → правила → Prometheus. "
               "Все параметры в конфигурации, артефакты сохранены."),
    (C_ACCENT, "Безучительский ансамбль эффективен для burst",
               "ROC-AUC 0.955, recall 0.97 без доступа к содержимому пакетов. "
               "Подход применим к зашифрованному трафику."),
    (C_ORANGE, "FPR требует калибровки — пути известны",
               "Повысить квантиль до 0.995; предфильтровать короткие потоки; "
               "добавить признак unique_dst_ports. Ожидаемое снижение FPR до 0.10–0.15."),
]
y = 1.35
for color, title, body in conclusions:
    add_rect(sl, 0.3, y, 12.73, 1.55, C_LIGHT)
    add_rect(sl, 0.3, y, 0.12, 1.55, color)
    add_text(sl, title, 0.55, y + 0.1, 12.1, 0.42, size=14, bold=True, color=color)
    add_text(sl, body, 0.55, y + 0.55, 12.1, 0.85, size=12.5, color=C_GRAY)
    y += 1.72

# Направления развития
add_rect(sl, 0.3, 6.3, 12.73, 0.9, C_DARK)
add_text(sl,
         "Направления:  сценарий port scan  ·  тест на реальном трафике  ·  "
         "адаптивный порог  ·  интеграция с SIEM  ·  расширение признаков (JA3, DNS)",
         0.5, 6.35, 12.3, 0.8, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Готово: {OUT}  ({OUT.stat().st_size // 1024} KB)  — 12 слайдов")
